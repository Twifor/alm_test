import collections
import collections.abc
from pptx import Presentation
import requests
import os
import time
import json
import sys
from agent.tools import Tool
from typing import Union, Dict

CWD = os.getcwd()  # path of current working directory
LIB_DIR = os.path.dirname(__file__)  # path of library
TEMPLATE_DIR = os.path.join(LIB_DIR, "templates")  # path of templates
CACHE_DIR = os.path.join(CWD, "cache")  # path of cache_dir
IMAGE_BED_PATTERN = (
    "https://source.unsplash.com/featured/?{}"  # url pattern for image bed
)
if not os.path.exists(CACHE_DIR):
    os.makedirs(CACHE_DIR)


def _return_timestamp():
    return str(time.time())


def runtime_update_docstring(new_docstring: str) -> callable:
    def decorator(func: callable) -> callable:
        func.__doc__ = new_docstring
        return func

    return decorator


ppt_template_names = []
all_files = os.listdir(TEMPLATE_DIR)
for file_name in all_files:
    if file_name.lower().endswith(".pptx"):
        ppt_template_names.append(file_name.split(".")[0])
updated_docstring_create_file = f"""create_file(theme:str) -> None: Create a pptx file with specific theme, available thems: {' / '.join(ppt_template_names)}."""
ppt_file = None


class CreateFileTool(Tool):
    def __init__(self):
        super().__init__()
        self.invoke_label = "CreateFile"

    @runtime_update_docstring(updated_docstring_create_file)
    def invoke(self, invoke_data) -> Union[str, int, bool, Dict]:
        theme = invoke_data
        global ppt_file
        ppt_file = Presentation(os.path.join(TEMPLATE_DIR, f"{theme}.pptx"))
        return "created a ppt file.", 0, False, {}

    def description(self) -> str:
        return "CreateFile(theme), create a pptx file with specific themes. Available themes: flat, green, orange, tech, wooden."


class GetImageTool(Tool):
    def __init__(self):
        super().__init__()
        self.invoke_label = "GetImage"

    def invoke(self, invoke_data) -> Union[str, int, bool, Dict]:
        keywords = invoke_data
        picture_url = IMAGE_BED_PATTERN.format(keywords)
        response = requests.get(picture_url)
        img_local_path = os.path.join(CACHE_DIR, f"{_return_timestamp()}.jpg")
        with open(img_local_path, "wb") as f:
            f.write(response.content)
        return img_local_path, 0, False, {}

    def description(self) -> str:
        return "GetImage(keywords), Get an image given comma separated keywords, return the image path."


class AddFirstPageTool(Tool):
    def __init__(self):
        super().__init__()
        self.invoke_label = "AddFirstPage"

    def invoke(self, invoke_data) -> Union[str, int, bool, Dict]:
        global ppt_file
        title, subtitle = invoke_data.split(",")
        slide = ppt_file.slides.add_slide(
            ppt_file.slide_layouts[0]
        )  # layout for first page (title and subtitle only)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title_shape.text = title
        subtitle_shape.text = subtitle

        return "added page", 0, False, {}

    def description(self) -> str:
        return "AddFirstPage(title, subtitle), Add the first page of ppt."


class AddTextPageTool(Tool):
    def __init__(self):
        super().__init__()
        self.invoke_label = "AddTextPage"

    def invoke(self, invoke_data) -> Union[str, int, bool, Dict]:
        title, bullet_items = invoke_data.split(",")
        global ppt_file
        slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[1])
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]

        title_shape.text = title

        tf = body_shape.text_frame

        bullet_items = bullet_items.split("[SPAN]")
        for bullet_item in bullet_items:
            bullet_item_strip = bullet_item.strip()
            p = tf.add_paragraph()
            p.text = bullet_item_strip
            p.level = 1

        return "added page", 0, False, {}

    def description(self) -> str:
        return "AddTextPage(title, bullet_items), Add text page (outline page is also applied). bullet_items should be string, for multiple bullet items, please use [SPAN] to separate them."


class AddTextImagePageTool(Tool):
    def __init__(self):
        super().__init__()
        self.invoke_label = "AddTextImagePage"

    def invoke(self, invoke_data) -> Union[str, int, bool, Dict]:
        title, bullet_items, image = invoke_data.split(",")
        global ppt_file
        slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[3])
        title_shape = slide.shapes.title
        title_shape.text = title

        body_shape = slide.placeholders[1]

        tf = body_shape.text_frame
        bullet_items = bullet_items.split("[SPAN]")
        for bullet_item in bullet_items:
            bullet_item_strip = bullet_item.strip()
            p = tf.add_paragraph()
            p.text = bullet_item_strip
            p.level = 1

        image_shape = slide.placeholders[2]
        slide.shapes.add_picture(
            image,
            image_shape.left,
            image_shape.top,
            image_shape.width,
            image_shape.height,
        )

        return "added page", 0, False, {}

    def description(self) -> str:
        return "AddTextImagePage(title, bullet_items, image), Add a text page with one image. (image should be a path). bullet_items should be string, for multiple bullet items, please use [SPAN] to separate them."


class SubmitFileTool(Tool):
    def __init__(self):
        super().__init__()
        self.invoke_label = "SubmitFile"

    def invoke(self, invoke_data="") -> Union[str, int, bool, Dict]:
        global ppt_file
        file_path = os.path.join(CACHE_DIR, f"{_return_timestamp()}.pptx")
        ppt_file.save(file_path)
        # retreival_url = upload_file(file_path)

        return f"submitted. view ppt at {file_path}", 0, False, {}

    def description(self) -> str:
        return "SubmitFile(title, bullet_items, image), When all steps done, YOU MUST use submit_file() to submit your ppt."
